from pptx import Presentation
from pptx.util import Inches


presentation = Presentation()


slide = presentation.slides.add_slide(presentation.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Stock Price Prediction"
subtitle.text = "Using Machine Learning to Forecast Stock Prices"


slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Introduction & Problem Statement"
content.text = (
    "The stock market is inherently volatile, and predicting stock prices "
    "is a challenging task influenced by numerous factors.\n\n"
    "The objective of this project is to develop a machine learning model "
    "to predict stock prices using historical data, enabling investors "
    "to make informed decisions."
)


slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Objectives"
content.text = (
    "1. Collect and preprocess historical stock data.\n"
    "2. Develop a machine learning model to predict stock prices.\n"
    "3. Evaluate the model's performance using metrics like Mean Squared Error (MSE).\n"
    "4. Visualize predictions and analyze insights.\n"
    "5. Provide recommendations for future work."
)


slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Methodology Overview"
content.text = (
    "The project methodology includes the following steps:\n"
    "1. Data Collection\n"
    "2. Data Preprocessing\n"
    "3. Feature Engineering\n"
    "4. Model Selection and Training\n"
    "5. Model Evaluation and Visualization"
)


slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Methodology - Data Collection"
content.text = (
    "1. Historical stock data was collected using the Yahoo Finance API.\n"
    "2. Example:\n"
    "   ```python\n"
    "   import yfinance as yf\n"
    "   data = yf.download('AAPL', start='2010-01-01', end='2023-01-01')\n"
    "   ```"
)


slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Methodology - Data Preprocessing"
content.text = (
    "1. Checked for missing values and handled them appropriately.\n"
    "2. Engineered new features like year, month, and day.\n"
    "3. Scaled data to ensure uniformity in feature ranges.\n"
    "4. Example:\n"
    "   ```python\n"
    "   data['Year'] = data.index.year\n"
    "   data['Month'] = data.index.month\n"
    "   ```"
)


slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Methodology - Model Selection and Training"
content.text = (
    "1. Model: Random Forest Regressor.\n"
    "2. Splitting data into training and testing sets.\n"
    "3. Example Code:\n"
    "   ```python\n"
    "   from sklearn.ensemble import RandomForestRegressor\n"
    "   model = RandomForestRegressor()\n"
    "   model.fit(X_train, y_train)\n"
    "   ```"
)


slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Results - Model Evaluation"
content.text = (
    "1. Evaluation Metrics:\n"
    "   - Mean Squared Error (MSE): Measures prediction error.\n"
    "   - R-squared: Proportion of variance explained.\n"
    "2. Example Code:\n"
    "   ```python\n"
    "   from sklearn.metrics import mean_squared_error, r2_score\n"
    "   mse = mean_squared_error(y_test, y_pred)\n"
    "   r2 = r2_score(y_test, y_pred)\n"
    "   ```"
)


slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Results - Visualization"
content.text = (
    "1. Plotted Actual vs Predicted Prices to evaluate trends.\n"
    "2. Example Code:\n"
    "   ```python\n"
    "   plt.plot(y_test.values, label='Actual')\n"
    "   plt.plot(y_pred, label='Predicted')\n"
    "   plt.legend()\n"
    "   plt.show()\n"
    "   ```"
)


slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Conclusion"
content.text = (
    "1. The model successfully captured stock price trends using historical data.\n"
    "2. Strengths:\n"
    "   - High accuracy for trend prediction.\n"
    "   - Feature importance analysis provides valuable insights.\n"
    "3. Limitations:\n"
    "   - External market factors were not considered.\n"
    "   - Sensitive to sudden market changes."
)


slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Future Work"
content.text = (
    "1. Implement advanced models like LSTMs for better time-series predictions.\n"
    "2. Integrate real-time external factors like news sentiment and economic data.\n"
    "3. Develop a user-friendly application for real-time stock price predictions.\n"
    "4. Enhance feature engineering with technical indicators (e.g., moving averages)."
)


presentation_path = "/mnt/data/Detailed_Project_Presentation.pptx"
presentation.save(presentation_path)
presentation_path

{
    
}
